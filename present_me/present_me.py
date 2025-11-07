import os
import time
import base64
import requests
import pandas as pd
from pathlib import Path
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Inches

# ---------- Config ----------
folder_path = os.path.expanduser("~/Desktop/present_me")
excel_file = os.path.join(folder_path, "data.xlsx")
template_file = os.path.join(folder_path, "template.html")
output_folder = os.path.join(folder_path, "output_images")
resource_images_folder = os.path.join(folder_path, "resource_images")

os.makedirs(output_folder, exist_ok=True)
os.makedirs(resource_images_folder, exist_ok=True)

# Optional logos & strip
logo_left = os.path.join(folder_path, "logo_left.png")
logo_right = os.path.join(folder_path, "logo_right.png")
bottom_strip = os.path.join(folder_path, "bottom_strip.png")

# Read template
with open(template_file, "r", encoding="utf-8") as f:
    template_html = f.read()

# HTML2Image Setup
hti = Html2Image(output_path=output_folder)
hti.custom_flags = [
    "--window-size=2500,2000",
    "--disable-web-security",
    "--allow-file-access-from-files",
    "--virtual-time-budget=20000"
]

# ---------- Networking: retry + caching ----------
download_cache_bytes = {}   # url -> bytes

def download_with_retry(url, retries=3, delay=5):
    headers = {"User-Agent": "Mozilla/5.0"}
    last_err = None
    for attempt in range(retries):
        try:
            r = requests.get(url, timeout=20, headers=headers)
            r.raise_for_status()
            return r.content
        except Exception as e:
            last_err = e
            print(f"⚠️ Attempt {attempt+1} failed for {url}: {e}")
            time.sleep(delay * (attempt + 1))
    print(f"❌ Failed after {retries} attempts for {url}: {last_err}")
    return None

def guess_mime_from_ext(path_or_url: str) -> str:
    lower = path_or_url.lower()
    if lower.endswith(".png"): return "image/png"
    if lower.endswith(".jpg") or lower.endswith(".jpeg"): return "image/jpeg"
    if lower.endswith(".gif"): return "image/gif"
    if lower.endswith(".webp"): return "image/webp"
    return "image/jpeg"

def bytes_to_data_uri(content: bytes, mime: str) -> str:
    b64 = base64.b64encode(content).decode("ascii")
    return f"data:{mime};base64,{b64}"

def file_to_data_uri(path: str) -> str:
    p = Path(path)
    if not p.exists():
        print(f"⚠️ Missing local image: {path}")
        return ""
    mime = guess_mime_from_ext(str(p))
    with open(p, "rb") as f:
        content = f.read()
    return bytes_to_data_uri(content, mime)

def url_to_data_uri(url: str, idx: int) -> str:
    url = (url or "").strip()
    if not url:
        return ""

    # Use cache (for speed)
    if url in download_cache_bytes:
        content = download_cache_bytes[url]
        mime = guess_mime_from_ext(url)
        return bytes_to_data_uri(content, mime)

    # Remote URL
    if url.startswith("http://") or url.startswith("https://"):
        content = download_with_retry(url, retries=4, delay=4)
        if not content:
            return ""

        # Store
        download_cache_bytes[url] = content
        mime = guess_mime_from_ext(url)
        return bytes_to_data_uri(content, mime)

    # Local path
    local_path = os.path.join(folder_path, url) if not os.path.isabs(url) else url
    if os.path.exists(local_path):
        return file_to_data_uri(local_path)

    print(f"⚠️ Invalid path: {url}")
    return ""

# ---------- HTML assembly ----------
def fill_template(template: str, row: pd.Series, idx: int) -> str:
    filled = template

    mapping = {
        "{{col2}}": row.get("col2", ""),
        "{{col3}}": row.get("col3", ""),
        "{{col4}}": row.get("col4", ""),
        "{{col5}}": row.get("col5", ""),
        "{{col6}}": row.get("col6", ""),
    }

    for placeholder, value in mapping.items():
        filled = filled.replace(placeholder, str(value if value is not None else ""))

    # Convert the main image into data URI
    img_data_uri = url_to_data_uri(str(row.get("col7", "")).strip(), idx)
    filled = filled.replace("{{col7}}", img_data_uri)

    # Convert logo paths into data URIs
    filled = filled.replace("{{logo_left}}", file_to_data_uri(logo_left))
    filled = filled.replace("{{logo_right}}", file_to_data_uri(logo_right))
    filled = filled.replace("{{bottom_strip}}", file_to_data_uri(bottom_strip))

    return filled

def render_image(row: pd.Series, idx: int, sheet_name: str) -> str:
    html_content = fill_template(template_html, row, idx)

    # Save debug HTML for troubleshooting
    debug_html_path = os.path.join(output_folder, f"{sheet_name}_page_{idx+1}.debug.html")
    with open(debug_html_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    output_filename = f"page_{idx+1}.png"
    output_path = os.path.join(output_folder, output_filename)

    hti.screenshot(html_str=html_content, save_as=output_filename)
    print(f"✅ Rendered page {idx+1} to {output_path}")
    return output_path

# ---------- PPT Assembly ----------
def add_simple_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(logo_left):
        slide.shapes.add_picture(logo_left, Inches(0.2), Inches(0.2), height=Inches(1))
    if os.path.exists(logo_right):
        slide.shapes.add_picture(logo_right, prs.slide_width - Inches(1.5), Inches(0.2), height=Inches(1))
    if os.path.exists(bottom_strip):
        slide.shapes.add_picture(bottom_strip, 0, prs.slide_height - Inches(1), width=prs.slide_width)
    return slide

def create_ppt(image_paths, ppt_file):
    prs = Presentation()
    add_simple_slide(prs)  # opening slide
    for img_path in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_simple_slide(prs)  # closing slide

    prs.save(ppt_file)
    print(f"✅ PPT saved: {ppt_file}")

# ---------- Main ----------
print("=== Robust Present-Me Generator ===")
xls = pd.ExcelFile(excel_file)

for sheet_name in xls.sheet_names:
    print(f"\nGenerating PPT for sheet: {sheet_name}")
    df_sheet = pd.read_excel(excel_file, sheet_name=sheet_name)

    image_paths = []
    for idx in range(len(df_sheet)):
        try:
            image_paths.append(render_image(df_sheet.iloc[idx], idx, sheet_name))
        except Exception as e:
            print(f"❌ Failed on row {idx+1}: {e}")

    ppt_file = os.path.join(folder_path, f"{sheet_name}.pptx")
    create_ppt(image_paths, ppt_file)
